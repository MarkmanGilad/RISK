#!/usr/bin/env python
"""EditorLM local extractor for one approved source."""

import argparse
import csv
import hashlib
import io
import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path, PurePosixPath


SUPPORTED = {
    ".html": "html", ".htm": "html", ".rtf": "rtf", ".csv": "csv", ".xlsx": "excel", ".xls": "excel",
    ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
}
HTML_BLOCKS = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "blockquote", "pre", "table"}


def utc_now():
    return datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")


def clean(value):
    return re.sub(r"\s+", " ", str(value)).strip()


def console_safe(value):
    """Return text that can be printed even in a legacy Windows console."""
    encoding = sys.stdout.encoding or "utf-8"
    return str(value).encode(encoding, errors="backslashreplace").decode(encoding)


def source_path(value):
    value = value.replace("\\", "/")
    path = PurePosixPath(value)
    if path.is_absolute() or ".." in path.parts or len(path.parts) < 2 or path.parts[0] != "sources":
        raise ValueError("--source must be a relative file path beneath sources/.")
    if path.suffix.lower() not in SUPPORTED:
        raise ValueError("This source format has no local extractor.")
    return path.as_posix()


def artifact_id(relative_path):
    stem = PurePosixPath(relative_path).stem
    slug = re.sub(r"[^A-Za-z0-9]+", "-", stem).strip("-").lower() or "source"
    digest = hashlib.sha256(relative_path.encode("utf-8")).hexdigest()[:12]
    return f"{slug}--{digest}"


def append_block(lines, entries, location, kind, text):
    if not text:
        return
    if lines:
        lines.append("")
    start = len(lines) + 1
    lines.extend(text.splitlines())
    entries.append({
        "sourceLocation": location,
        "kind": kind,
        "parsedTextLineStart": start,
        "parsedTextLineEnd": len(lines),
    })


def with_links(element):
    text = clean(element.get_text(" ", strip=True))
    links = []
    for anchor in element.find_all("a", href=True):
        label = clean(anchor.get_text(" ", strip=True)) or anchor["href"]
        links.append(f"[{label}]({anchor['href']})")
    return f"{text} — Links: {', '.join(links)}" if links else text


def extract_html(path):
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(path.read_bytes(), "html.parser")
    for tag in soup.find_all(("script", "style", "noscript", "template")):
        tag.decompose()
    root = soup.body or soup
    lines, entries, ordinals = [], [], {}
    if soup.title:
        append_block(lines, entries, "title[1]", "title", f"# {clean(soup.title.get_text(' ', strip=True))}")
    for element in root.find_all(list(HTML_BLOCKS)):
        if element.find_parent(HTML_BLOCKS):
            continue
        tag = element.name.lower()
        ordinals[tag] = ordinals.get(tag, 0) + 1
        location = f"{tag}[{ordinals[tag]}]"
        if tag == "table":
            for row_number, row in enumerate(element.find_all("tr"), start=1):
                values = [with_links(cell) for cell in row.find_all(("th", "td"), recursive=False)]
                append_block(lines, entries, f"{location}/row[{row_number}]", "table-row", " | ".join(values))
            continue
        text = with_links(element)
        if tag.startswith("h") and len(tag) == 2 and tag[1].isdigit():
            text = f"{'#' * int(tag[1])} {text}"
        elif tag == "li":
            text = f"- {text}"
        elif tag == "blockquote":
            text = f"> {text}"
        append_block(lines, entries, location, tag, text)
    if not lines:
        append_block(lines, entries, "document[1]", "document", clean(root.get_text(" ", strip=True)))
    return "\n".join(lines).strip() + "\n", entries


def extract_rtf(path):
    raw = path.read_bytes().decode("latin-1")
    code_page_match = re.search(r"\\ansicpg(\d+)", raw[:8192])
    code_page = f"cp{code_page_match.group(1)}" if code_page_match else "cp1252"
    try:
        "".encode(code_page)
    except LookupError:
        code_page = "cp1252"
    ignored_destinations = {
        "fonttbl", "colortbl", "stylesheet", "info", "pict", "object", "header", "footer",
        "footnote", "annotation", "field", "fldinst", "fldrslt", "datastore", "themedata",
    }
    output, states = [], [False]
    index, unicode_fallback_count, fallback_remaining = 0, 1, 0
    while index < len(raw):
        character = raw[index]
        if character == "{":
            states.append(states[-1])
            index += 1
            continue
        if character == "}":
            if len(states) > 1:
                states.pop()
            index += 1
            continue
        if character != "\\":
            if fallback_remaining > 0:
                fallback_remaining -= 1
            elif not states[-1]:
                output.append(character)
            index += 1
            continue

        index += 1
        if index >= len(raw):
            break
        control = raw[index]
        if control in "\\{}":
            if fallback_remaining > 0:
                fallback_remaining -= 1
            elif not states[-1]:
                output.append(control)
            index += 1
            continue
        if control == "'" and index + 2 < len(raw):
            try:
                decoded = bytes([int(raw[index + 1:index + 3], 16)]).decode(code_page)
                if fallback_remaining > 0:
                    fallback_remaining -= 1
                elif not states[-1]:
                    output.append(decoded)
            except (LookupError, UnicodeDecodeError, ValueError):
                pass
            index += 3
            continue
        if control == "~":
            if not states[-1]:
                output.append(" ")
            index += 1
            continue
        if control in "-_":
            if control == "_" and not states[-1]:
                output.append("-")
            index += 1
            continue
        if control == "*":
            states[-1] = True
            index += 1
            continue

        match = re.match(r"([A-Za-z]+)(-?\d+)? ?", raw[index:])
        if not match:
            index += 1
            continue
        word, number = match.group(1), match.group(2)
        index += len(match.group(0))
        if word in ignored_destinations:
            states[-1] = True
        elif word in {"par", "line"} and not states[-1]:
            output.append("\n")
        elif word == "tab" and not states[-1]:
            output.append("\t")
        elif word == "uc" and number is not None:
            unicode_fallback_count = max(0, int(number))
        elif word == "u" and number is not None:
            value = int(number)
            if value < 0:
                value += 65536
            if not states[-1]:
                output.append(chr(value))
            fallback_remaining = unicode_fallback_count

    paragraphs = [clean(paragraph) for paragraph in "".join(output).splitlines()]
    lines, entries = [], []
    for paragraph_number, paragraph in enumerate(paragraphs, start=1):
        append_block(lines, entries, f"paragraph[{paragraph_number}]", "rtf-paragraph", paragraph)
    if not lines:
        raise ValueError("No readable text was extracted from this RTF source.")
    return "\n".join(lines).strip() + "\n", entries


def extract_csv(path):
    raw = path.read_bytes()
    decoded = None
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            decoded = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            pass
    if decoded is None:
        raise ValueError("CSV encoding could not be read.")
    dialect = csv.excel
    try:
        dialect = csv.Sniffer().sniff(decoded[:4096], delimiters=",;\t|")
    except csv.Error:
        pass
    rows = list(csv.reader(io.StringIO(decoded), dialect))
    lines, entries = [], []
    for row_number, row in enumerate(rows, start=1):
        append_block(lines, entries, f"row[{row_number}]", "csv-row", " | ".join(clean(value) for value in row))
    return "\n".join(lines).strip() + "\n", entries


def excel_value(value):
    return "" if value is None else clean(value)


def column_name(index):
    value = ""
    while index:
        index, remainder = divmod(index - 1, 26)
        value = chr(65 + remainder) + value
    return value


def extract_xlsx(path):
    from openpyxl import load_workbook
    workbook = load_workbook(path, read_only=True, data_only=False)
    lines, entries = [], []
    for worksheet in workbook.worksheets:
        append_block(lines, entries, f"sheet[{worksheet.title}]", "worksheet", f"## Sheet: {worksheet.title}")
        for row_number, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
            values = [excel_value(value) for value in row]
            append_block(lines, entries, f"{worksheet.title}!{column_name(1)}{row_number}", "worksheet-row", " | ".join(values))
            for column_number, value in enumerate(values, start=1):
                if value:
                    entries.append({
                        "sourceLocation": f"{worksheet.title}!{column_name(column_number)}{row_number}",
                        "kind": "cell",
                        "parsedTextLineStart": len(lines),
                        "parsedTextLineEnd": len(lines),
                    })
    return "\n".join(lines).strip() + "\n", entries


def extract_xls(path):
    import xlrd
    workbook = xlrd.open_workbook(path)
    lines, entries = [], []
    for worksheet in workbook.sheets():
        append_block(lines, entries, f"sheet[{worksheet.name}]", "worksheet", f"## Sheet: {worksheet.name}")
        for row_index in range(worksheet.nrows):
            values = [excel_value(worksheet.cell_value(row_index, column)) for column in range(worksheet.ncols)]
            row_number = row_index + 1
            append_block(lines, entries, f"{worksheet.name}!{column_name(1)}{row_number}", "worksheet-row", " | ".join(values))
            for column_number, value in enumerate(values, start=1):
                if value:
                    entries.append({
                        "sourceLocation": f"{worksheet.name}!{column_name(column_number)}{row_number}",
                        "kind": "cell",
                        "parsedTextLineStart": len(lines),
                        "parsedTextLineEnd": len(lines),
                    })
    return "\n".join(lines).strip() + "\n", entries


def extract_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path, strict=False)
    lines, entries = [], []
    for page_number, page in enumerate(reader.pages, start=1):
        layout_text = page.extract_text(extraction_mode="layout") or ""
        plain_text = page.extract_text() or ""
        text = plain_text if len(clean(plain_text)) > len(clean(layout_text)) else layout_text
        append_block(lines, entries, f"page[{page_number}]", "pdf-page", clean(text))
    if not lines:
        raise ValueError("No selectable text was extracted from this PDF. It may require the later OCR or image-vision workflow.")
    return "\n".join(lines).strip() + "\n", entries


def extract_docx(path):
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    document = Document(path)
    lines, entries, paragraph_number, table_number = [], [], 0, 0
    for block in document.iter_inner_content():
        if isinstance(block, Paragraph):
            paragraph_number += 1
            text = clean(block.text)
            style = block.style.name if block.style else ""
            match = re.match(r"Heading ([1-6])", style)
            if match:
                text = f"{'#' * int(match.group(1))} {text}"
            append_block(lines, entries, f"paragraph[{paragraph_number}]", "paragraph", text)
        elif isinstance(block, Table):
            table_number += 1
            for row_number, row in enumerate(block.rows, start=1):
                values = [clean(cell.text) for cell in row.cells]
                append_block(lines, entries, f"table[{table_number}]/row[{row_number}]", "table-row", " | ".join(values))
    return "\n".join(lines).strip() + "\n", entries


def extract_pptx(path):
    from pptx import Presentation
    presentation = Presentation(path)
    lines, entries = [], []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        append_block(lines, entries, f"slide[{slide_number}]", "slide", f"# Slide {slide_number}")
        for shape_number, shape in enumerate(slide.shapes, start=1):
            location = f"slide[{slide_number}]/shape[{shape_number}]"
            if getattr(shape, "has_table", False):
                for row_number, row in enumerate(shape.table.rows, start=1):
                    values = [clean(cell.text) for cell in row.cells]
                    append_block(lines, entries, f"{location}/row[{row_number}]", "table-row", " | ".join(values))
            elif getattr(shape, "has_text_frame", False):
                append_block(lines, entries, location, "shape-text", clean(shape.text))
        notes = slide.notes_slide
        if notes and notes.notes_text_frame:
            append_block(lines, entries, f"slide[{slide_number}]/notes", "speaker-notes", clean(notes.notes_text_frame.text))
    return "\n".join(lines).strip() + "\n", entries


def extract(path, format_name):
    if format_name == "html":
        return extract_html(path)
    if format_name == "rtf":
        return extract_rtf(path)
    if format_name == "csv":
        return extract_csv(path)
    if format_name == "excel":
        return extract_xlsx(path) if path.suffix.lower() == ".xlsx" else extract_xls(path)
    if format_name == "pdf":
        return extract_pdf(path)
    if format_name == "docx":
        return extract_docx(path)
    if format_name == "pptx":
        return extract_pptx(path)
    raise ValueError("No local extractor is available for this source format.")


def read_registry(path):
    try:
        parsed = json.loads(path.read_text(encoding="utf-8"))
        if parsed.get("schemaVersion") == 1 and isinstance(parsed.get("sources"), dict):
            return parsed
    except (FileNotFoundError, json.JSONDecodeError):
        pass
    return {"schemaVersion": 1, "updatedAt": "1970-01-01T00:00:00Z", "sources": {}}


def save_state(project_root, registry, entry):
    state = project_root / ".editorlm" / "state"
    state.mkdir(parents=True, exist_ok=True)
    registry["updatedAt"] = utc_now()
    (state / "source-registry.json").write_text(json.dumps(registry, indent=2) + "\n", encoding="utf-8")
    with (state / "processing-log.jsonl").open("a", encoding="utf-8") as log:
        log.write(json.dumps(entry) + "\n")


def main():
    parser = argparse.ArgumentParser(description="Extract one approved local source for EditorLM.")
    parser.add_argument("--project-root", default=".", help="EditorLM project root (defaults to current folder).")
    parser.add_argument("--source", required=True, help="Relative source path beneath sources/.")
    args = parser.parse_args()
    project_root = Path(args.project_root).resolve()
    relative_path = source_path(args.source)
    original = (project_root / Path(*PurePosixPath(relative_path).parts)).resolve()
    try:
        original.relative_to((project_root / "sources").resolve())
    except ValueError as error:
        raise SystemExit("--source must resolve inside sources/.") from error
    if not original.is_file():
        raise SystemExit(f"Source does not exist: {relative_path}")
    format_name = SUPPORTED[original.suffix.lower()]
    registry = read_registry(project_root / ".editorlm" / "state" / "source-registry.json")
    now = utc_now()
    try:
        text, entries = extract(original, format_name)
        identifier = artifact_id(relative_path)
        parsed_relative = f".editorlm/processed/{identifier}/parsed-text.txt"
        map_relative = f".editorlm/processed/{identifier}/location-map.json"
        parsed_path = project_root / Path(*PurePosixPath(parsed_relative).parts)
        map_path = project_root / Path(*PurePosixPath(map_relative).parts)
        parsed_path.parent.mkdir(parents=True, exist_ok=True)
        parsed_path.write_text(text, encoding="utf-8")
        map_path.write_text(json.dumps({
            "schemaVersion": 1,
            "format": format_name,
            "sourcePath": relative_path,
            "parsedTextPath": parsed_relative,
            "entries": entries,
        }, indent=2) + "\n", encoding="utf-8")
        stats = original.stat()
        registry["sources"][relative_path] = {
            "format": format_name,
            "lastProcessed": {
                "sourceModifiedAtMs": stats.st_mtime_ns / 1_000_000,
                "sourceSizeBytes": stats.st_size,
                "processedAt": now,
                "processorVersion": 1,
                "kind": f"{format_name}-extracted",
                "parsedTextPath": parsed_relative,
                "locationMapPath": map_relative,
            },
            "lastAttempt": {"at": now, "outcome": "processed"},
        }
        save_state(project_root, registry, {
            "at": now,
            "operation": f"{format_name}-extraction",
            "relativePath": relative_path,
            "outcome": "processed",
            "sourceSizeBytes": stats.st_size,
        })
        print(f"Extracted {console_safe(relative_path)} to {parsed_relative}")
    except Exception as error:
        existing = registry["sources"].get(relative_path, {})
        registry["sources"][relative_path] = {
            "format": format_name,
            "lastProcessed": existing.get("lastProcessed"),
            "lastAttempt": {"at": now, "outcome": "failed", "error": str(error)},
        }
        save_state(project_root, registry, {
            "at": now,
            "operation": f"{format_name}-extraction",
            "relativePath": relative_path,
            "outcome": "failed",
            "error": str(error),
        })
        raise SystemExit(f"EditorLM extraction failed: {error}") from error


if __name__ == "__main__":
    main()
